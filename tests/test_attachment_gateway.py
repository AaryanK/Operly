import asyncio,io,json,tempfile,unittest,zipfile
from pathlib import Path
from PIL import Image
from reportlab.pdfgen import canvas
from docx import Document
from pptx import Presentation
from openpyxl import Workbook,load_workbook
from packages.business_brain.attachments.detector import detect_type,DetectionError
from packages.business_brain.attachments.formatter import requested_format,operation,split_discord_text
from packages.business_brain.attachments.limits import AttachmentLimits
from packages.business_brain.attachments.models import AttachmentInput,AttachmentBundle,ParsedAttachment
from packages.business_brain.attachments.models import GeneratedOutput,OutputFile
from packages.business_brain.attachments.multimodal_processor import MultimodalProcessor,SYSTEM_PROMPT
from packages.business_brain.attachments.outputs import generate_output,safe_cell
from packages.business_brain.attachments.parsers import parse_attachment
from packages.business_brain.attachments.archive_processor import process_archive
from packages.business_brain.attachments.privacy import infer_sensitivity

def image_bytes(color="red"):
    out=io.BytesIO();Image.new("RGB",(20,10),color).save(out,"PNG");return out.getvalue()
def pdf_bytes(text=""):
    out=io.BytesIO();c=canvas.Canvas(out)
    if text:c.drawString(50,750,text)
    c.showPage();c.save();return out.getvalue()
def docx_bytes():
    out=io.BytesIO();d=Document();d.add_paragraph("Traveler Alice");t=d.add_table(rows=2,cols=2);t.cell(0,0).text="Name";t.cell(1,0).text="Alice";d.save(out);return out.getvalue()
def pptx_bytes():
    out=io.BytesIO();p=Presentation();s=p.slides.add_slide(p.slide_layouts[1]);s.shapes.title.text="Travel";s.placeholders[1].text="Group package";p.save(out);return out.getvalue()
def xlsx_bytes(formula=False):
    out=io.BytesIO();w=Workbook();s=w.active;s.title="Travelers";s.append(["Name","Due"]);s.append(["Alice","=NOW()" if formula else "2026-01-01"]);w.save(out);return out.getvalue()
def item(i,name,data,declared=None):
    x=AttachmentInput(i,name,declared,len(data),data);x.detected_content_type=detect_type(name,data,declared);return x
class FakeClient:
    def __init__(self):self.calls=[]
    async def chat(self,messages,tools):self.calls.append(messages);return {"content":f"analysis-{len(self.calls)}"}

class DetectionParsingTests(unittest.TestCase):
    def test_one_image(self):self.assertEqual(parse_attachment(item(1,"a.png",image_bytes())).metadata["width"],20)
    def test_image_ordering(self):
        rows=[parse_attachment(item(i,f"{i}.png",image_bytes())) for i in range(1,6)];self.assertEqual([x.index for x in rows],[1,2,3,4,5])
    def test_unreadable_image(self):
        with self.assertRaises((DetectionError,ValueError)):parse_attachment(item(1,"a.png",b"\x89PNG\r\n\x1a\nBAD"))
    def test_text_pdf(self):self.assertIn("Hello PDF",parse_attachment(item(1,"a.pdf",pdf_bytes("Hello PDF"))).extracted_text)
    def test_scanned_pdf(self):
        row=parse_attachment(item(1,"scan.pdf",pdf_bytes()),10);self.assertEqual(row.metadata["pages"],1);self.assertTrue(row.images or row.warnings)
    def test_docx(self):self.assertIn("Alice",parse_attachment(item(1,"a.docx",docx_bytes())).extracted_text)
    def test_pptx(self):self.assertIn("Group package",parse_attachment(item(1,"a.pptx",pptx_bytes())).extracted_text)
    def test_xlsx_formula_inert(self):self.assertIn("=NOW()",parse_attachment(item(1,"a.xlsx",xlsx_bytes(True))).extracted_text)
    def test_csv_json_and_source(self):
        csvrow=parse_attachment(item(1,"a.csv",b"name,due\nA,now"));self.assertEqual(csvrow.tables[0][0][0],"name")
        self.assertEqual(parse_attachment(item(2,"a.json",b'{"ok":true}')).category,"text")
        self.assertIn("print",parse_attachment(item(3,"a.py",b"print('inert')")).extracted_text)
    def test_executable_rejected(self):
        with self.assertRaises(DetectionError):detect_type("bad.exe",b"MZpayload")
    def test_office_mismatch(self):
        with self.assertRaises(DetectionError):detect_type("fake.docx",b"plain text")
    def test_actual_signature_wins_mime(self):self.assertEqual(detect_type("x.bin",pdf_bytes(),"image/png"),"application/pdf")
    def test_malformed_document(self):
        with self.assertRaises(DetectionError):detect_type("bad.xlsx",b"PK\x03\x04broken")
    def test_sensitive_inference(self):self.assertEqual(infer_sensitivity("passport number A12345678"),"highly_sensitive")
    def test_prompt_injection_remains_data(self):self.assertIn("untrusted",SYSTEM_PROMPT.lower())
class ArchiveTests(unittest.TestCase):
    def make_zip(self,entries):
        out=io.BytesIO()
        with zipfile.ZipFile(out,"w") as z:
            for name,data in entries:z.writestr(name,data)
        return out.getvalue()
    def test_supported_zip(self):
        data=self.make_zip([("one.txt",b"hello"),("table.csv",b"a,b\n1,2")]);rows=process_archive(item(1,"a.zip",data),AttachmentLimits());self.assertEqual(len(rows),2)
    def test_path_traversal_not_written(self):
        data=self.make_zip([("../../secret.txt",b"safe")]);rows=process_archive(item(1,"a.zip",data),AttachmentLimits());self.assertTrue(rows[0].filename.endswith("secret.txt"));self.assertNotIn("..",rows[0].filename)
    def test_entry_limit(self):
        data=self.make_zip([(f"{i}.txt",b"x") for i in range(3)]);limits=AttachmentLimits(max_archive_files=2)
        with self.assertRaises(ValueError):process_archive(item(1,"a.zip",data),limits)
    def test_expanded_limit(self):
        data=self.make_zip([("big.txt",b"x"*1000)]);limits=AttachmentLimits(max_archive_expanded_bytes=100)
        with self.assertRaises(ValueError):process_archive(item(1,"a.zip",data),limits)
class OutputTests(unittest.TestCase):
    def setUp(self):self.parsed=[ParsedAttachment(1,"a.txt","text","text/plain","hello")];self.analyses={1:"=danger result"}
    def test_formula_safety(self):self.assertEqual(safe_cell("=1+1"),"'=1+1")
    def test_json_csv_docx_pdf_xlsx_outputs(self):
        with tempfile.TemporaryDirectory() as folder:
            for fmt in ("json","csv","xlsx","docx","pdf","markdown","txt"):
                files=generate_output(fmt,self.parsed,self.analyses,"summary",folder);self.assertTrue(files[0].path.is_file());self.assertGreater(files[0].size_bytes,0)
            wb=load_workbook(Path(folder)/"operly-attachment-result.xlsx",data_only=False);self.assertTrue(str(wb.active["F2"].value).startswith("'="))
    def test_discord_chunks(self):
        chunks=split_discord_text("word "*1000,190);self.assertTrue(all(len(x)<=190 for x in chunks));self.assertGreater(len(chunks),1)
    def test_intent_output(self):self.assertEqual(requested_format("return Excel"),"xlsx");self.assertEqual(operation("compare these files"),"compare")
    def test_temporary_cleanup_contract(self):
        with tempfile.TemporaryDirectory() as folder:path=generate_output("json",self.parsed,self.analyses,"summary",folder)[0].path;self.assertTrue(path.exists())
        self.assertFalse(path.exists())
class ProcessorTests(unittest.IsolatedAsyncioTestCase):
    async def test_five_images_and_order(self):
        client=FakeClient();items=[AttachmentInput(i,f"{i}.png","image/png",len(image_bytes()),image_bytes()) for i in range(1,6)];result=await MultimodalProcessor(client).process(AttachmentBundle("Describe each",items));self.assertEqual(result.accepted,[f"{i}.png" for i in range(1,6)]);self.assertEqual(len(client.calls),5)
    async def test_explicit_comparison(self):
        client=FakeClient();items=[AttachmentInput(i,f"{i}.txt","text/plain",1,b"x") for i in (1,2)];result=await MultimodalProcessor(client).process(AttachmentBundle("Compare these",items));self.assertEqual(result.operation_summary,"compare");self.assertEqual(len(client.calls),3)
    async def test_requested_xlsx(self):
        client=FakeClient();x=AttachmentInput(1,"a.csv","text/csv",3,b"a,b")
        with tempfile.TemporaryDirectory() as folder:
            result=await MultimodalProcessor(client).process(AttachmentBundle("return Excel",[x]),folder);self.assertEqual(result.files[0].filename,"operly-attachment-result.xlsx")
    async def test_acceptance_seven_files_compare_to_excel(self):
        client=FakeClient();png=image_bytes();attachments=[AttachmentInput(i,f"traveler-{i}.png","image/png",len(png),png) for i in range(1,6)]
        pdf=pdf_bytes("Traveler manifest");sheet=xlsx_bytes();attachments += [AttachmentInput(6,"booking.pdf","application/pdf",len(pdf),pdf),AttachmentInput(7,"travelers.xlsx","application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",len(sheet),sheet)]
        with tempfile.TemporaryDirectory() as folder:
            result=await MultimodalProcessor(client).process(AttachmentBundle("Extract traveler information, compare it with the spreadsheet, identify inconsistencies and return Excel",attachments),folder)
            self.assertEqual(len(result.accepted),7);self.assertEqual(result.operation_summary,"compare");self.assertEqual(result.files[0].filename,"operly-attachment-result.xlsx");self.assertEqual(len(client.calls),8)
    async def test_limits(self):
        limits=AttachmentLimits(max_attachments=1,max_total_bytes=3,max_attachment_bytes=3);p=MultimodalProcessor(FakeClient(),limits)
        with self.assertRaises(ValueError):await p.process(AttachmentBundle("x",[AttachmentInput(1,"a.txt",None,1,b"x"),AttachmentInput(2,"b.txt",None,1,b"x")]))
        with self.assertRaises(ValueError):await p.process(AttachmentBundle("x",[AttachmentInput(1,"a.txt",None,4,b"xxxx")]))

class DiscordGatewayTests(unittest.IsolatedAsyncioTestCase):
    async def test_generated_file_response_and_cleanup(self):
        from unittest.mock import AsyncMock,patch
        from packages.connectors.discord import bot_shared
        class Processor:
            limits=AttachmentLimits()
            async def process(self,bundle,temp_dir):
                path=Path(temp_dir)/"result.json";path.write_text("{}")
                self.path=path
                return GeneratedOutput("Complete",[OutputFile(path,"result.json","application/json",2)],accepted=["input.txt"])
        class Attachment:
            filename="input.txt";content_type="text/plain";size=5
            async def read(self):return b"hello"
        class Progress:
            def __init__(self):self.edits=[];self.author=type("A",(),{"id":99,"display_name":"OPERLY"})();self.guild=None;self.channel=channel;self.id=999
            async def edit(self,**kwargs):self.edits.append(kwargs)
        class Channel:
            id=42
            def __init__(self):self.files=[];self.messages=[]
            async def send(self,*args,**kwargs):
                self.messages.append(args[0] if args else "")
                if kwargs.get("file"):
                    self.files.append(kwargs["file"].filename);kwargs["file"].close()
        channel=Channel();progress=Progress()
        message=type("M",(),{})();message.id=10;message.attachments=[Attachment()];message.author=type("U",(),{"id":7})();message.guild=None;message.channel=channel;message.reply=AsyncMock(return_value=progress)
        processor=Processor()
        with patch.object(bot_shared,"attachment_processor",processor),patch.object(bot_shared,"attachment_already_processed",AsyncMock(return_value=False)),patch.object(bot_shared,"audit_attachments",AsyncMock()),patch.object(bot_shared,"store_message",AsyncMock()):
            await bot_shared.process_discord_attachments(message,"tenant-1","return JSON")
        self.assertEqual(channel.files,["result.json"]);self.assertFalse(processor.path.exists());self.assertTrue(progress.edits)
if __name__=="__main__":unittest.main()
