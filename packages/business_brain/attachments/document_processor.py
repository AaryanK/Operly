import base64,io,json
from .models import AttachmentInput,ParsedAttachment
from .privacy import infer_sensitivity

def _bounded(value,limit=200_000):
    value=value or "";return value[:limit],(["extracted text was truncated"] if len(value)>limit else [])
def _safe_image(blob):
    try:
        from PIL import Image
        with Image.open(io.BytesIO(blob)) as image:
            if image.width*image.height>36_000_000:return None
            image.load();out=io.BytesIO();image.convert("RGB").save(out,"PNG")
        return base64.b64encode(out.getvalue()).decode("ascii")
    except Exception:return None
def parse_pdf(item,max_pages):
    try:from pypdf import PdfReader
    except ImportError as exc:raise RuntimeError("PDF parser dependency is unavailable") from exc
    reader=PdfReader(io.BytesIO(item.content_bytes),strict=False)
    if len(reader.pages)>max_pages:raise ValueError(f"PDF exceeds {max_pages} page limit")
    texts=[];warnings=[]
    for i,page in enumerate(reader.pages):
        try:texts.append(f"[Page {i+1}]\n{page.extract_text() or ''}")
        except Exception:texts.append(f"[Page {i+1}]\n[unreadable]");warnings.append(f"page {i+1} could not be read")
    text="\n\n".join(texts)
    if len(text.strip())<40:
        try:
            import fitz
            doc=fitz.open(stream=item.content_bytes,filetype="pdf")
            images=[]
            for page in doc:
                pix=page.get_pixmap(matrix=fitz.Matrix(1.5,1.5),alpha=False);images.append(__import__("base64").b64encode(pix.tobytes("png")).decode("ascii"))
            warnings.append("PDF contained little embedded text; rendered pages were supplied as images")
        except Exception:images=[];warnings.append("PDF appears scanned and page rendering is unavailable")
    else:images=[]
    text,extra=_bounded(text);warnings+=extra
    return ParsedAttachment(item.index,item.filename,"document",item.detected_content_type,text,images=images,metadata={"pages":len(reader.pages)},warnings=warnings,sensitivity=infer_sensitivity(text))
def parse_docx(item):
    try:from docx import Document
    except ImportError as exc:raise RuntimeError("DOCX parser dependency is unavailable") from exc
    from docx.text.paragraph import Paragraph
    from docx.table import Table
    from docx.oxml.text.paragraph import CT_P
    from docx.oxml.table import CT_Tbl
    doc=Document(io.BytesIO(item.content_bytes));parts=[];tables=[];images=[]
    for child in doc.element.body.iterchildren():
        if isinstance(child,CT_P):
            value=Paragraph(child,doc).text
            if value.strip():parts.append(value)
        elif isinstance(child,CT_Tbl):
            table=Table(child,doc);rows=[[cell.text for cell in row.cells] for row in table.rows];tables.append(rows);parts.append("\n".join(" | ".join(r) for r in rows))
    for part in doc.part.related_parts.values():
        if getattr(part,"content_type","").startswith("image/"):
            value=_safe_image(part.blob)
            if value:images.append(value)
    text,w=_bounded("\n".join(parts));return ParsedAttachment(item.index,item.filename,"document",item.detected_content_type,text,tables=tables,images=images,metadata={"paragraphs":len(doc.paragraphs),"tables":len(tables)},warnings=w,sensitivity=infer_sensitivity(text))
def parse_pptx(item):
    try:from pptx import Presentation
    except ImportError as exc:raise RuntimeError("PPTX parser dependency is unavailable") from exc
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs=Presentation(io.BytesIO(item.content_bytes));parts=[];images=[]
    for n,slide in enumerate(prs.slides,1):
        texts=[]
        for shape in slide.shapes:
            if hasattr(shape,"text") and shape.text.strip():texts.append(shape.text)
            if shape.shape_type==MSO_SHAPE_TYPE.PICTURE:
                value=_safe_image(shape.image.blob)
                if value:images.append(value)
        if slide.has_notes_slide:
            notes=[x.text for x in slide.notes_slide.notes_text_frame.paragraphs if x.text.strip()]
            if notes:texts.append("Notes: "+" ".join(notes))
        parts.append(f"[Slide {n}]\n"+"\n".join(texts))
    text,w=_bounded("\n\n".join(parts));return ParsedAttachment(item.index,item.filename,"presentation",item.detected_content_type,text,images=images,metadata={"slides":len(prs.slides)},warnings=w,sensitivity=infer_sensitivity(text))
def parse_xlsx(item):
    try:from openpyxl import load_workbook
    except ImportError as exc:raise RuntimeError("XLSX parser dependency is unavailable") from exc
    wb=load_workbook(io.BytesIO(item.content_bytes),read_only=True,data_only=False,keep_links=False);tables=[];parts=[]
    for ws in wb.worksheets[:30]:
        rows=[]
        for r_idx,row in enumerate(ws.iter_rows(max_row=5000,max_col=100),1):
            vals=[]
            for cell in row:
                value=cell.value
                vals.append("" if value is None else str(value)[:2000])
            if any(vals):rows.append(vals)
        tables.append(rows);parts.append(f"[Sheet: {ws.title}]\n"+"\n".join("\t".join(r) for r in rows))
    text,w=_bounded("\n\n".join(parts));return ParsedAttachment(item.index,item.filename,"spreadsheet",item.detected_content_type,text,tables=tables,metadata={"sheets":wb.sheetnames},warnings=w,sensitivity=infer_sensitivity(text))
def parse_odf(item):
    try:
        from odf.opendocument import load
        from odf import text,table
    except ImportError as exc:raise RuntimeError("OpenDocument parser dependency is unavailable") from exc
    doc=load(io.BytesIO(item.content_bytes));values=[]
    for node in doc.getElementsByType(text.P):values.append("".join(str(x) for x in node.childNodes))
    extracted,w=_bounded("\n".join(values));return ParsedAttachment(item.index,item.filename,"document",item.detected_content_type,extracted,warnings=w,sensitivity=infer_sensitivity(extracted))
